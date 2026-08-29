# coding:UTF-8
# Build the sensor comparison deck from the measured numbers.
#
#   python make_pptx.py                 -> SENSOR_COMPARISON.pptx at the repo root
#   python make_pptx.py --out other.pptx
#
# The numbers live here as literals rather than being re-derived from the CSVs.
# The captures were made on a bench that no longer exists in that state, so
# recomputing them is not reproducible anyway, and a deck that silently changes
# when a stray CSV is edited is worse than one that has to be updated on purpose.
# SENSOR_COMPARISON.md is the source of record; this file mirrors it.
#
# Fonts are the ones a Korean Windows install actually has. The web version uses
# IBM Plex, which would silently fall back to something arbitrary in PowerPoint.

from __future__ import annotations

import argparse
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# The palette validated for colour-vision deficiency in the web report.
HWT = RGBColor(0xB8, 0x5C, 0x00)
SCL = RGBColor(0x00, 0x90, 0xC4)
INK = RGBColor(0x0F, 0x1A, 0x21)
INK2 = RGBColor(0x3A, 0x4C, 0x57)
MUTED = RGBColor(0x65, 0x79, 0x85)
RULE = RGBColor(0xD5, 0xE0, 0xE7)
GROUND = RGBColor(0xEE, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FLAG = RGBColor(0x9A, 0x34, 0x12)

SANS = "맑은 고딕"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = W - 2 * MARGIN


def textbox(slide, left, top, width, height, text, size=18, color=INK,
            bold=False, font=SANS, align=PP_ALIGN.LEFT, space_after=6,
            line_spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rule(slide, top, left=MARGIN, width=BODY_W, color=RULE, height=Pt(1.25)):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def new_slide(prs, eyebrow, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    rule(slide, Inches(0.62))
    textbox(slide, MARGIN, Inches(0.76), BODY_W, Inches(0.3), eyebrow,
            size=12, color=MUTED, bold=True)
    textbox(slide, MARGIN, Inches(1.06), BODY_W, Inches(0.7), title,
            size=30, color=INK, bold=True)
    top = Inches(1.84)
    if subtitle:
        # Size the box to the wrapped text and start the body below it. A fixed
        # height here is what put the subtitle on top of the first table.
        sub_w = Inches(11.4)
        per_line = 74                      # characters that fit at 15 pt in sub_w
        lines = sum(max(1, -(-len(seg) // per_line))
                    for seg in subtitle.split("\n"))
        box_h = Inches(0.27) * lines + Inches(0.04)
        textbox(slide, MARGIN, top, sub_w, box_h, subtitle, size=15, color=INK2)
        top = top + box_h + Inches(0.3)
    return slide, top


def table(slide, left, top, rows, widths, header=True, size=13,
          col_colors=None, width_total=None, row_h=Pt(30)):
    """rows[0] is the header. col_colors maps a column index to a body colour."""
    n_rows, n_cols = len(rows), len(rows[0])
    total = width_total or sum(widths)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, total,
                                   row_h * n_rows)
    tbl = shape.table
    for i, w in enumerate(widths):
        tbl.columns[i].width = w
    for r, row in enumerate(rows):
        tbl.rows[r].height = row_h
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = GROUND if r == 0 and header else WHITE
            tf = cell.text_frame
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(size)
            run.font.name = MONO if (c > 0 and r > 0) else SANS
            run.font.bold = (r == 0 and header)
            if r == 0 and header:
                run.font.color.rgb = MUTED
            elif col_colors and c in col_colors:
                run.font.color.rgb = col_colors[c]
            else:
                run.font.color.rgb = INK
    return shape


def stat_row(slide, top, items, height=Inches(1.15)):
    """items: (label, value, unit, note, colour)"""
    from pptx.enum.shapes import MSO_SHAPE
    gap = Inches(0.22)
    width = Emu(int((BODY_W - gap * (len(items) - 1)) / len(items)))
    for i, (label, value, unit, note, colour) in enumerate(items):
        left = MARGIN + Emu(int(i * (width + gap)))
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = GROUND
        card.line.color.rgb = RULE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        textbox(slide, left + Inches(0.18), top + Inches(0.13),
                width - Inches(0.3), Inches(0.24), label, size=11,
                color=MUTED, bold=True)
        vbox = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.4),
                                        width - Inches(0.3), Inches(0.42))
        tf = vbox.text_frame
        tf.margin_left = tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = value
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.name = MONO
        run.font.color.rgb = colour
        if unit:
            u = para.add_run()
            u.text = " " + unit
            u.font.size = Pt(13)
            u.font.name = SANS
            u.font.color.rgb = MUTED
        if note:
            textbox(slide, left + Inches(0.18), top + Inches(0.83),
                    width - Inches(0.3), Inches(0.24), note, size=11, color=INK2)


def log_axis(chart):
    """python-pptx exposes no log scale, and this data spans two decades."""
    from pptx.oxml.ns import qn
    from lxml import etree
    for ax in chart._chartSpace.iter(qn("c:valAx")):
        scaling = ax.find(qn("c:scaling"))
        if scaling is None:
            continue
        base = etree.Element(qn("c:logBase"))
        base.set("val", "10")
        scaling.insert(0, base)


def averaging_chart(slide, left, top, width, height):
    data = CategoryChartData()
    data.categories = ["1", "2", "5", "10", "25", "50", "100"]
    data.add_series("HWT9037 실측", (0.0005, 0.0005, 0.0005, 0.0005, 0.0004, 0.0004, 0.0004))
    data.add_series("SCL3300 실측", (0.0098, 0.0073, 0.0033, 0.0023, 0.0014, 0.0009, 0.0006))
    data.add_series("백색잡음이라면 (SCL)", (0.0098, 0.0069, 0.0044, 0.0031, 0.0020, 0.0014, 0.0010))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top,
                                   width, height, data)
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = SANS
    chart.font.size = Pt(11)
    chart.font.name = SANS

    colours = [HWT, SCL, MUTED]
    for series, colour in zip(chart.plots[0].series, colours):
        series.format.line.color.rgb = colour
        series.format.line.width = Pt(2.25)
        series.smooth = False
    dashed = chart.plots[0].series[2]
    dashed.format.line.width = Pt(1.5)

    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.name = MONO
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = MONO
    log_axis(chart)
    return chart


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------

def build(path: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # --- 1. title ---------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, GROUND)
    textbox(slide, MARGIN, Inches(2.15), BODY_W, Inches(0.3),
            "VERTICRANE · 경사 센서 벤치 비교", size=13, color=MUTED, bold=True)
    textbox(slide, MARGIN, Inches(2.62), BODY_W, Inches(1.2),
            "HWT9037  vs  SCL3300", size=54, color=INK, bold=True)
    textbox(slide, MARGIN, Inches(3.95), Inches(8.2), Inches(1.0),
            "원시 잡음 수치는 두 센서의 순위를 뒤집어 놓는다.\n"
            "그 수치가 왜 틀렸고, 대역폭을 맞추면 무엇이 남는가.",
            size=17, color=INK2)
    rule(slide, Inches(5.35), width=Inches(2.2), color=SCL, height=Pt(3))
    textbox(slide, MARGIN, Inches(5.62), Inches(9.0), Inches(1.1),
            "2026-08-29  ·  Raspberry Pi Zero 2 W  ·  Modbus RTU + SPI0 CE1\n"
            "잡음 150 s / 25 Hz / 3,750 샘플  ·  주파수응답 3분  ·  스케일 12–14° 3쌍",
            size=12, color=MUTED, font=MONO, line_spacing=1.5)

    # --- 2. verdict -------------------------------------------------------
    slide, top = new_slide(prs, "결론", "재려는 것에 따라 답이 갈린다")
    table(slide, MARGIN, top, [
        ["재려는 것", "판정", "근거"],
        ["1 % 기울기 같은 정적 각도", "동등", "오차는 센서가 아니라 설치가 지배"],
        ["지속적인 흔들림", "동등", "최악 조합조차 잡음의 252배"],
        ["삐걱거림 같은 짧고 작은 외란", "SCL3300", "HWT는 약 10배 큰 외란을 요구"],
    ], [Inches(4.6), Inches(2.0), Inches(5.29)], size=15, row_h=Pt(46),
        col_colors={1: SCL})
    textbox(slide, MARGIN, Inches(5.05), BODY_W, Inches(1.4),
            "마지막 줄이 결정적일 수 있다. HWT9037은 2–3 Hz에서 꺾이고 작은 신호는 양자화에 막히므로,\n"
            "크레인이 삐걱거리는 중에도 \"안정\"이라고 답할 수 있다.",
            size=16, color=INK)

    # --- 3. the misleading first answer ------------------------------------
    slide, top = new_slide(prs, "먼저 나온 답", "원시 수치는 HWT9037의 압승이라고 말한다",
                           "정지 상태 150초, 각자의 각도 출력을 그대로 놓고 잰 표준편차.")
    stat_row(slide, top, [
        ("HWT9037 ANGX σ", "0.0005", "°", "peak-to-peak 0.0030°", HWT),
        ("SCL3300 ANGX σ", "0.0097", "°", "peak-to-peak 0.0879°", SCL),
        ("겉보기 차이", "19", "×", "HWT9037이 조용해 보인다", INK),
    ])
    textbox(slide, MARGIN, Inches(4.5), BODY_W, Inches(1.2),
            "여기서 멈추면 결론은 하나뿐이다.\n"
            "그런데 같은 데이터에 이 결론과 양립할 수 없는 숫자가 둘 있다.",
            size=17, color=INK)

    # --- 4. rebuttal one ---------------------------------------------------
    slide, top = new_slide(prs, "반증 하나", "HWT9037의 잡음이 자기 눈금보다 325배 작다",
                           "정지한 센서의 |acc|는 중력 그대로 1 g여야 하고, 그 흔들림은 순수한 측정 잡음이다.")
    table(slide, MARGIN, top, [
        ["항목", "HWT9037-485", "SCL3300"],
        ["분해능", "0.488 mg/LSB", "0.083 mg/LSB"],
        ["|acc| σ", "0.0015 mg", "0.1508 mg"],
        ["150초 peak-to-peak", "0.008 mg", "2.428 mg"],
        ["σ ÷ 자기 LSB", "1/325", "1.8"],
    ], [Inches(4.6), Inches(3.6), Inches(3.69)], size=14,
        col_colors={1: HWT, 2: SCL})
    textbox(slide, MARGIN, Inches(5.45), BODY_W, Inches(1.2),
            "ADC 한 칸을 못 넘는 잡음은 필터 없이 존재할 수 없다. 150초 동안 움직인 전체 폭이 자기 LSB의 1/61이라는 것은,\n"
            "그 시간 내내 출력이 사실상 상수였다는 뜻이다.",
            size=15, color=INK)

    # --- 5. rebuttal two ---------------------------------------------------
    slide, top = new_slide(prs, "반증 둘", "평균을 내도 나아지지 않는 출력",
                           "백색잡음은 N개를 평균 내면 1/√N로 줄어든다. 이미 평활된 출력은 꿈쩍하지 않는다.")
    averaging_chart(slide, MARGIN, top, Inches(7.4), Inches(4.2))
    textbox(slide, Inches(8.5), top + Inches(0.3), Inches(4.1), Inches(3.6),
            "HWT9037은 100개를 평균 내도\n수평을 유지한다.\n\n"
            "연속 샘플이 서로 독립이 아니라는 뜻이고,\n"
            "그 값은 잡음이 아니라 필터의 바닥이다.\n\n"
            "SCL3300은 1/√N를 따라 내려가\nN=100에서 0.0006°로 만난다.\n\n"
            "가로축은 평균 구간, 세로축은 각도 σ(도), 로그.",
            size=13, color=INK2)

    # --- 6. the cost -------------------------------------------------------
    slide, top = new_slide(prs, "대가", "필터는 진짜로 일어난 일도 지운다",
                           "측정 중 책상을 건드렸다. 두 센서가 같은 순간에 5σ 이벤트를 잡았으니 실재한 사건이다.")
    stat_row(slide, top, [
        ("충격 중 최대 편차 · HWT9037", "0.98", "mg", "자기 LSB의 2칸", HWT),
        ("충격 중 최대 편차 · SCL3300", "21.4", "mg", "258 LSB", SCL),
        ("감쇠비", "22", "×", "주파수 때문이지 진폭이 아니다", INK),
    ])
    textbox(slide, MARGIN, Inches(4.45), BODY_W, Inches(1.6),
            "각도 출력으로 보면 HWT9037은 0.004° 범위에만 머물렀고 SCL3300은 0.35° 움직였다.\n\n"
            "이 장비의 판정 로직이 하는 일은 \"크레인이 안정되었는가\"를 묻는 것이다.\n"
            "짧은 흔들림을 지워버리는 출력은 흔들리는 중에도 \"안정\"이라고 답할 수 있다.",
            size=15, color=INK)

    # --- 7. frequency response ---------------------------------------------
    slide, top = new_slide(prs, "흔들어서 측정", "필터의 코너는 2–3 Hz다",
                           "잡음 바닥의 40배(100–200 mg)로 3분을 흔들고 4초 창으로 훑어 주파수별로 분류했다.")
    table(slide, MARGIN, top, [
        ["대역", "창", "HWT9037", "SCL3300", "비", "1차 fc"],
        ["0.00 – 0.35 Hz", "24", "203.7 mg", "209.8 mg", "1.03", "—"],
        ["0.35 – 0.70 Hz", "5", "133.6 mg", "132.3 mg", "0.99", "—"],
        ["0.70 – 1.40 Hz", "4", "113.8 mg", "119.1 mg", "1.05", "3.22 Hz"],
        ["1.40 – 2.80 Hz", "3", "103.5 mg", "122.9 mg", "1.19", "3.09 Hz"],
        ["2.80 – 5.60 Hz", "9", "19.1 mg", "42.6 mg", "2.24", "1.98 Hz"],
    ], [Inches(2.7), Inches(0.8), Inches(2.3), Inches(2.3), Inches(1.5), Inches(2.29)],
        size=13, col_colors={2: HWT, 3: SCL})
    textbox(slide, MARGIN, Inches(5.55), BODY_W, Inches(1.2),
            "HWT9037은 1.4 Hz까지 SCL3300을 5 % 이내로 따라온다 — \"느린 것도 못 본다\"는 인상은 틀렸다.\n"
            "충격의 22배 감쇠는 진폭이 아니라 주파수 때문이었다. 순간 충격 에너지는 이 곡선의 오른쪽 바깥에 있다.",
            size=15, color=INK)

    # --- 8. detection margin ------------------------------------------------
    slide, top = new_slide(prs, "검출 여유", "지속적인 흔들림은 두 센서 모두 여유롭게 잡는다",
                           "가진을 각도로 환산하고 각 센서의 정지 잡음으로 나눈 값.")
    table(slide, MARGIN, top, [
        ["가진", "HWT9037", "잡음 대비", "SCL3300", "잡음 대비"],
        ["느림 (0–0.35 Hz)", "11.67° rms", "23,000배", "12.02° rms", "1,240배"],
        ["중간 (0.7–1.4 Hz)", "6.52° rms", "13,000배", "6.82° rms", "703배"],
        ["빠름 (2.8–5.6 Hz)", "1.09° rms", "2,180배", "2.44° rms", "252배"],
    ], [Inches(3.0), Inches(2.3), Inches(2.0), Inches(2.3), Inches(2.29)],
        size=14, col_colors={1: HWT, 3: SCL})
    textbox(slide, MARGIN, Inches(4.9), BODY_W, Inches(1.4),
            "가장 불리한 조합조차 잡음의 252배다. 이 요구사항만 놓고 보면 두 센서는 교체 가능하다.\n"
            "갈리는 것은 크기가 아니라 지속 시간이다.",
            size=16, color=INK)

    # --- 9. resolution ------------------------------------------------------
    slide, top = new_slide(prs, "눈금과 재현성", "분해능과 정밀도는 다른 것이다",
                           "사양이 아니라 150초 기록에서 실제로 나타난 값들의 최소 간격.")
    table(slide, MARGIN, top, [
        ["필드", "실측 눈금", "부품 사양", "150초간 고유값"],
        ["HWT9037 AccX/Y", "1.0 mg", "0.488 mg", "2개"],
        ["HWT9037 AccZ", "—", "0.488 mg", "1개"],
        ["SCL3300 Acc", "0.08 mg", "0.083 mg", "20–25개"],
        ["HWT9037 Ang", "0.001°", "0.001°", "4–5개"],
        ["SCL3300 Ang", "0.0054°", "0.00549°", "17–18개"],
    ], [Inches(3.6), Inches(2.6), Inches(2.6), Inches(3.09)], size=13)
    textbox(slide, MARGIN, Inches(5.5), BODY_W, Inches(1.3),
            "드라이버가 HWT9037 가속도를 1 mg로 반올림해 부품 분해능의 절반을 버리고 있다. AccZ는 150초 동안 값이 하나뿐이었다.\n"
            "각도 눈금은 HWT9037이 5.4배 곱지만, 그 가속도계 1 LSB는 각도로 0.028°다 — 0.001°는 분해가 아니라 보간이다.",
            size=14, color=INK)

    # --- 10. bandwidth ------------------------------------------------------
    slide, top = new_slide(prs, "대역폭", "SCL3300은 느리게 읽는 것이 아니라 평균 내야 한다",
                           "내부 LPF가 10 Hz이고 출력 레이트 제어가 없다. 띄엄띄엄 읽으면 0.5–10 Hz 잡음이 접혀 들어온다.")
    table(slide, MARGIN, top, [
        ["출력 레이트", "솎아내기 σ", "평균 σ", "이득"],
        ["5 Hz", "0.00993°", "0.00325°", "3.1배"],
        ["1 Hz", "0.00854°", "0.00137°", "6.2배"],
        ["0.5 Hz", "0.00880°", "0.00087°", "10.2배"],
        ["0.2 Hz", "0.00810°", "0.00062°", "13.1배"],
    ], [Inches(3.0), Inches(3.0), Inches(3.0), Inches(2.89)], size=14,
        col_colors={2: SCL})
    textbox(slide, MARGIN, Inches(5.15), BODY_W, Inches(1.4),
            "솎아내기는 출력을 아무리 낮춰도 σ가 0.009°에서 내려오지 않는다. N개 평균은 그 자체가 안티앨리어싱 필터다.\n\n"
            "25 Hz는 HWT9037에는 과잉이지만 SCL3300에는 최적이다 — LPF 10 Hz에 Nyquist 12.5 Hz로 맞고, 이웃 샘플이 독립이다.",
            size=15, color=INK)

    # --- 11. one percent ----------------------------------------------------
    slide, top = new_slide(prs, "실사용 판단 (1)", "1 % 기울기 — 둘 다 문제없다",
                           "1 % 기울기는 0.5729°다. 오차 항목을 같은 단위에 놓으면 승부가 센서에 있지 않다.")
    table(slide, MARGIN, top, [
        ["항목", "각도", "% 기울기"],
        ["HWT9037 정밀도 (25 Hz)", "0.0005°", "0.0009 %"],
        ["SCL3300 정밀도 (0.5 Hz)", "0.0014°", "0.0024 %"],
        ["SCL3300 정밀도 (25 Hz)", "0.0097°", "0.0169 %"],
        ["HWT9037 드라이버 반올림 후", "0.0573°", "0.1000 %"],
        ["자세 재현성 (설치 오차)", "0.2500°", "0.4363 %"],
    ], [Inches(5.6), Inches(3.1), Inches(3.19)], size=14,
        col_colors={2: INK})
    textbox(slide, MARGIN, Inches(5.5), BODY_W, Inches(1.2),
            "예산을 지배하는 것은 설치다. 자세 재현성 0.25°가 그 자체로 0.44 % 기울기이며, 가장 나쁜 센서 항목보다 4배 크다.\n"
            "같은 노력을 장착 정밀도와 영점 관리에 쓰는 편이 훨씬 크게 남는다.",
            size=15, color=INK)

    # --- 12. creaking -------------------------------------------------------
    slide, top = new_slide(prs, "실사용 판단 (2)", "삐걱거림 — 여기서는 SCL3300이다",
                           "삐걱거림은 짧고, 빠르고, 작다. HWT9037이 불리한 세 조건이 한꺼번에 걸린다.")
    stat_row(slide, top, [
        ("3–5 Hz 최소 검출 · HWT9037", "0.257", "°", "4.48 mg · 0.448 % 기울기", HWT),
        ("3–5 Hz 최소 검출 · SCL3300", "0.026", "°", "0.45 mg · 0.045 % 기울기", SCL),
        ("요구 외란 크기 차이", "10", "×", "HWT9037이 더 큰 것을 요구", FLAG),
    ])
    table(slide, MARGIN, Inches(4.35), [
        ["책상 충격이 더 작았다면", "HWT9037 카운트", "SCL3300 카운트"],
        ["실제 (21.4 mg)", "1", "258"],
        ["1/5", "0.2", "52"],
        ["1/20", "0", "13"],
    ], [Inches(5.6), Inches(3.1), Inches(3.19)], size=14,
        col_colors={1: HWT, 2: SCL})
    textbox(slide, MARGIN, Inches(6.35), BODY_W, Inches(0.8),
            "그 충격의 1/5 크기 외란은 HWT9037에게 완전히 보이지 않고, SCL3300에게는 52카운트로 또렷하다.",
            size=16, color=INK, bold=True)

    # --- 13. accuracy and scale ---------------------------------------------
    slide, top = new_slide(prs, "정확도와 스케일", "절대 정확도는 HWT9037, 스케일은 무승부",
                           "중력은 눈금이 필요 없는 유일한 기준이다. 스케일은 12–14° 회전 세 쌍으로 쟀다.")
    table(slide, MARGIN, top, [
        ["항목", "HWT9037-485", "SCL3300"],
        ["|acc| 평균", "1.00102 g", "0.99586 g"],
        ["1 g 대비 오차", "+0.10 %", "−0.41 %"],
        ["150초 드리프트", "<0.0002 °/분", "<0.0004 °/분"],
        ["스케일 비 (3쌍 평균)", "기준", "+0.6 % ± 1.4 %"],
    ], [Inches(4.6), Inches(3.6), Inches(3.69)], size=14,
        col_colors={1: HWT, 2: SCL})
    textbox(slide, MARGIN, Inches(5.45), BODY_W, Inches(1.3),
            "각도는 축들의 비에서 나오므로 세 축 공통 감도 오차는 상쇄된다. 그리고 감도 오차는 잡음과 달리 교정으로 없앨 수 있다.\n"
            "스케일은 세 쌍이 +2.15 / −0.70 / +0.33 %로 갈렸는데, 그 편차 자체가 측정 한계다 — 센서가 아니라 자세 재현성에서 온다.",
            size=14, color=INK)

    # --- 14. shared bus -----------------------------------------------------
    slide, top = new_slide(prs, "통합", "공유 SPI 버스는 영향이 없다",
                           "e-paper 패널이 CE0, SCL3300이 CE1. 패널이 전체 리프레시를 도는 동안 100 Hz로 읽었다.")
    table(slide, MARGIN, top, [
        ["항목", "SCL3300 단독", "e-paper 리프레시 동시"],
        ["샘플", "500 / 5 s", "2,000 / 20 s"],
        ["실측 레이트", "100.0 Hz", "100.0 Hz"],
        ["최악 루프 지연", "10.1 ms", "11.5 ms"],
        ["CRC 오류", "0", "0"],
    ], [Inches(4.6), Inches(3.6), Inches(3.69)], size=14, col_colors={2: SCL})
    textbox(slide, MARGIN, Inches(5.45), BODY_W, Inches(1.3),
            "간섭은 1.5 ms뿐이다. 1.4초 리프레시 중 실제로 버스를 점유하는 구간이 10 ms 남짓이고 나머지는 BUSY 폴링이기 때문이다.\n"
            "경합이 없는 이유는 셋이다 — 패널이 MISO를 쓰지 않고, 둘 다 SPI 모드 0이며, 커널이 전송을 직렬화한다.",
            size=14, color=INK)

    # --- 15. decision -------------------------------------------------------
    slide, top = new_slide(prs, "남은 결정", "대체인가 병행인가",
                           "비교는 끝났다. 갈림길은 SCL3300에 자이로가 없다는 점에서 갈린다.")
    table(slide, MARGIN, top, [
        ["", "대체", "병행"],
        ["변경 범위", "recorder.py 샘플 소스", "레코드 포맷 또는 파일 분리"],
        ["걸리는 곳", "128바이트 레코드의 자이로 필드", "기존 .dat 파일과의 호환"],
        ["확인 필요", "stability.py가 자이로를 쓰는지", "전력·CPU"],
    ], [Inches(2.6), Inches(4.6), Inches(4.69)], size=14, row_h=Pt(44))
    textbox(slide, MARGIN, Inches(5.15), BODY_W, Inches(1.5),
            "어느 쪽이든 삐걱거림 판정이 결정 기준에 들어가야 한다.\n"
            "잡음만 보면 HWT9037이 매력적이지만, 그 잡음이 낮은 이유와 짧은 흔들림을 못 보는 이유가 같은 것이다.",
            size=16, color=INK)
    rule(slide, Inches(6.6))
    textbox(slide, MARGIN, Inches(6.75), BODY_W, Inches(0.4),
            "SENSOR_COMPARISON.md · dev/sensor_compare.py · scl3300.py",
            size=11, color=MUTED, font=MONO)

    prs.save(path)
    print("{} 슬라이드 → {}".format(len(prs.slides._sldIdLst), path))


def main() -> int:
    root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ap = argparse.ArgumentParser(description="센서 비교 발표자료 생성")
    ap.add_argument("--out", default=os.path.join(root, "SENSOR_COMPARISON.pptx"))
    args = ap.parse_args()
    build(args.out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
